# Office documents processor
from pathlib import Path
from typing import List
import docx
from openpyxl import load_workbook
from pptx import Presentation
from .base import BaseProcessor, ProcessorResult

class OfficeProcessor(BaseProcessor):
    SUPPORTED_TYPES = {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
        "application/msword": "doc",
        "application/vnd.ms-excel": "xls",
    }
    
    def can_process(self, mime_type: str) -> bool:
        return mime_type in self.SUPPORTED_TYPES
    
    def process(self, file_path: Path, **kwargs) -> List[ProcessorResult]:
        mime_type = kwargs.get("mime_type", "")
        doc_type = self.SUPPORTED_TYPES.get(mime_type)
        
        if doc_type == "docx":
            return self._process_docx(file_path)
        elif doc_type == "xlsx":
            return self._process_xlsx(file_path)
        elif doc_type == "pptx":
            return self._process_pptx(file_path)
        else:
            return [ProcessorResult(content_text="Unsupported Office format")]
    
    def _process_docx(self, file_path: Path) -> List[ProcessorResult]:
        doc = docx.Document(file_path)
        text = "
".join([para.text for para in doc.paragraphs])
        return [ProcessorResult(content_text=text)]
    
    def _process_xlsx(self, file_path: Path) -> List[ProcessorResult]:
        wb = load_workbook(file_path, read_only=True)
        text_parts = []
        
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                text_parts.append(" | ".join(str(cell) if cell else "" for cell in row))
        
        return [ProcessorResult(content_text="
".join(text_parts))]
    
    def _process_pptx(self, file_path: Path) -> List[ProcessorResult]:
        prs = Presentation(file_path)
        text_parts = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        
        return [ProcessorResult(content_text="
".join(text_parts))]
